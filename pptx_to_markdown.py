"""Convert PowerPoint (.pptx) files to Markdown format.

Supports headings, bullet lists, tables, and basic inline formatting.
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER  # type: ignore
except ImportError:  # pragma: no cover - clear message if dependency missing
    Presentation = None  # type: ignore
    MSO_SHAPE_TYPE = None  # type: ignore
    PP_PLACEHOLDER = None  # type: ignore


def _runs_to_markdown_text(runs) -> str:
    """
    Convert a sequence of pptx paragraph runs to Markdown inline formatting.
    Supports bold, italic, and underline.

    Uses a state machine to avoid emitting separate markers for each run,
    mirroring the approach in docx_to_markdown.py.
    """
    result: list[str] = []
    cur_bold = False
    cur_italic = False
    cur_underline = False

    for r in runs:
        text = r.text or ""
        if not text:
            continue

        target_bold = bool(r.font.bold)
        target_italic = bool(r.font.italic)
        target_underline = bool(r.font.underline)

        # Close markers that are no longer active
        if cur_underline and not target_underline:
            result.append("</u>")
            cur_underline = False
        if cur_italic and not target_italic:
            result.append("_")
            cur_italic = False
        if cur_bold and not target_bold:
            result.append("**")
            cur_bold = False

        # Open newly active markers
        if target_bold and not cur_bold:
            result.append("**")
            cur_bold = True
        if target_italic and not cur_italic:
            result.append("_")
            cur_italic = True
        if target_underline and not cur_underline:
            result.append("<u>")
            cur_underline = True

        result.append(text)

    # Close any markers still open at the end of the paragraph
    if cur_underline:
        result.append("</u>")
    if cur_italic:
        result.append("_")
    if cur_bold:
        result.append("**")

    return "".join(result)


def _table_to_markdown(table) -> list[str]:
    """Convert a pptx table to a list of Markdown table lines."""
    rows = list(table.rows)
    if not rows:
        return []

    lines: list[str] = []
    header_cells = [cell.text.strip().replace("\n", " ") for cell in rows[0].cells]
    if any(header_cells):
        lines.append("| " + " | ".join(header_cells) + " |")
        lines.append("| " + " | ".join("---" for _ in header_cells) + " |")
        for row in rows[1:]:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")
    else:
        for row in rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            lines.append(" | ".join(cells))

    return lines


def _placeholder_type(shape):
    """Return the placeholder type of *shape*, or ``None`` if not a placeholder."""
    if not shape.is_placeholder:
        return None
    return shape.placeholder_format.type


def pptx_to_markdown(input_path: Path) -> str:
    """
    Convert a .pptx file to a Markdown string.

    Notes:
    - Slide titles become H1 headings.
    - Subtitles become H2 headings.
    - Body text is rendered as a bullet list, indented by paragraph level.
    - Tables become Markdown tables.
    - Images are emitted as ``![name](name)`` placeholders.
    - Speaker notes are appended as blockquotes.
    - Slides are separated by horizontal rules (``---``).
    """
    if Presentation is None:
        raise RuntimeError(
            "Missing dependency: python-pptx.\n"
            "Install it with: pip install python-pptx"
        )

    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")

    prs = Presentation(str(input_path))

    _title_types = frozenset(
        [
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
            PP_PLACEHOLDER.VERTICAL_TITLE,
        ]
    )
    _subtitle_types = frozenset([PP_PLACEHOLDER.SUBTITLE])

    md_lines: list[str] = []

    for slide_idx, slide in enumerate(prs.slides):
        # Separate slides with a horizontal rule (skip before the first slide)
        if slide_idx > 0:
            md_lines.append("---")
            md_lines.append("")

        for shape in slide.shapes:
            ph_type = _placeholder_type(shape)

            if shape.has_text_frame:
                if ph_type in _title_types:
                    # Slide title → H1
                    all_runs = [
                        run
                        for para in shape.text_frame.paragraphs
                        for run in para.runs
                    ]
                    text = _runs_to_markdown_text(all_runs).strip()
                    if not text:
                        text = shape.text_frame.text.strip()
                    if text:
                        md_lines.append(f"# {text}")
                        md_lines.append("")

                elif ph_type in _subtitle_types:
                    # Subtitle placeholder → H2
                    all_runs = [
                        run
                        for para in shape.text_frame.paragraphs
                        for run in para.runs
                    ]
                    text = _runs_to_markdown_text(all_runs).strip()
                    if not text:
                        text = shape.text_frame.text.strip()
                    if text:
                        md_lines.append(f"## {text}")
                        md_lines.append("")

                else:
                    # Body placeholder or free text box → bullet list
                    has_content = any(
                        para.text.strip()
                        for para in shape.text_frame.paragraphs
                    )
                    if not has_content:
                        continue

                    for para in shape.text_frame.paragraphs:
                        text = _runs_to_markdown_text(para.runs).strip()
                        if not text:
                            text = para.text.strip()
                        if not text:
                            md_lines.append("")
                            continue
                        level = para.level or 0
                        indent = "  " * level
                        md_lines.append(f"{indent}- {text}")

                    md_lines.append("")

            elif shape.has_table:
                table_lines = _table_to_markdown(shape.table)
                md_lines.extend(table_lines)
                md_lines.append("")

            elif (
                MSO_SHAPE_TYPE is not None
                and shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            ):
                alt = shape.name or "image"
                md_lines.append(f"![{alt}]({alt})")
                md_lines.append("")

        # Speaker notes as blockquote
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip() if notes_tf else ""
            if notes_text:
                md_lines.append(f"> **Notes:** {notes_text.replace('\n', ' ')}")
                md_lines.append("")

    # Remove trailing blank lines
    while md_lines and not md_lines[-1].strip():
        md_lines.pop()

    return "\n".join(md_lines)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Convert a .pptx file to Markdown."
    )
    parser.add_argument("input", help="Path to input .pptx file")
    parser.add_argument(
        "output",
        nargs="?",
        help="Path to output .md file (default: same name with .md extension)",
    )

    args = parser.parse_args(argv)

    input_path = Path(args.input).expanduser().resolve()
    if args.output:
        output_path = Path(args.output).expanduser().resolve()
    else:
        output_path = input_path.with_suffix(".md")

    try:
        md = pptx_to_markdown(input_path)
    except Exception as exc:  # pragma: no cover - CLI error path
        sys.stderr.write(f"Error: {exc}\n")
        return 1

    output_path.write_text(md, encoding="utf-8")
    print(f"Wrote Markdown to: {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
