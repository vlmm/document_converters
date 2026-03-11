import io
import sys
import os
import tempfile
import unittest

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from pptx_to_markdown import pptx_to_markdown, _runs_to_markdown_text, _table_to_markdown


def _make_pptx(*slide_builders) -> Path:
    """Helper: build a Presentation using callables, save to a temp file, return path."""
    prs = Presentation()
    for builder in slide_builders:
        builder(prs)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    tmp_dir = Path(tempfile.mkdtemp())
    tmp_path = tmp_dir / "test.pptx"
    tmp_path.write_bytes(buf.getvalue())
    return tmp_path


class TestRunsToMarkdownText(unittest.TestCase):
    """Unit tests for the _runs_to_markdown_text helper."""

    def _make_run(self, text, bold=None, italic=None, underline=None):
        """Return a simple mock run object."""

        class _Font:
            pass

        class _Run:
            pass

        run = _Run()
        run.text = text
        run.font = _Font()
        run.font.bold = bold
        run.font.italic = italic
        run.font.underline = underline
        return run

    def test_plain_text(self):
        runs = [self._make_run("Hello World")]
        self.assertEqual(_runs_to_markdown_text(runs), "Hello World")

    def test_bold(self):
        runs = [self._make_run("Bold", bold=True)]
        self.assertEqual(_runs_to_markdown_text(runs), "**Bold**")

    def test_italic(self):
        runs = [self._make_run("Italic", italic=True)]
        self.assertEqual(_runs_to_markdown_text(runs), "_Italic_")

    def test_underline(self):
        runs = [self._make_run("Under", underline=True)]
        self.assertEqual(_runs_to_markdown_text(runs), "<u>Under</u>")

    def test_mixed_formatting(self):
        runs = [
            self._make_run("Normal "),
            self._make_run("Bold", bold=True),
            self._make_run(" End"),
        ]
        self.assertEqual(_runs_to_markdown_text(runs), "Normal **Bold** End")

    def test_empty_runs(self):
        runs = [self._make_run("")]
        self.assertEqual(_runs_to_markdown_text(runs), "")


class TestTableToMarkdown(unittest.TestCase):
    """Unit tests for the _table_to_markdown helper."""

    def _make_table(self, rows_data):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        nrows = len(rows_data)
        ncols = max(len(r) for r in rows_data)
        table = slide.shapes.add_table(
            nrows, ncols, Inches(1), Inches(1), Inches(6), Inches(3)
        ).table
        for r_idx, row in enumerate(rows_data):
            for c_idx, cell_text in enumerate(row):
                table.cell(r_idx, c_idx).text = cell_text
        return table

    def test_table_with_header(self):
        table = self._make_table([["H1", "H2"], ["R1", "R2"]])
        lines = _table_to_markdown(table)
        self.assertIn("| H1 | H2 |", lines[0])
        self.assertIn("---", lines[1])
        self.assertIn("R1", lines[2])

    def test_table_empty(self):
        # Empty table should produce no lines (but python-pptx requires at least 1 row)
        table = self._make_table([["", ""]])
        lines = _table_to_markdown(table)
        # All-empty header → fallback plain text rows
        self.assertIsInstance(lines, list)


class TestPptxToMarkdown(unittest.TestCase):
    """Integration tests for pptx_to_markdown."""

    def test_title_slide(self):
        """Title and subtitle should become H1 and H2."""

        def build(prs):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "My Title"
            slide.placeholders[1].text = "My Subtitle"

        path = _make_pptx(build)
        md = pptx_to_markdown(path)
        self.assertIn("# My Title", md)
        self.assertIn("## My Subtitle", md)

    def test_bullet_list_levels(self):
        """Body text paragraphs at different levels become indented bullets."""

        def build(prs):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Bullets"
            body = slide.placeholders[1]
            body.text_frame.text = "Level 0"
            p1 = body.text_frame.add_paragraph()
            p1.text = "Level 1"
            p1.level = 1

        path = _make_pptx(build)
        md = pptx_to_markdown(path)
        self.assertIn("- Level 0", md)
        self.assertIn("  - Level 1", md)

    def test_table_conversion(self):
        """Tables should become Markdown tables."""

        def build(prs):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            table = slide.shapes.add_table(
                2, 2, Inches(1), Inches(1), Inches(4), Inches(2)
            ).table
            table.cell(0, 0).text = "Col A"
            table.cell(0, 1).text = "Col B"
            table.cell(1, 0).text = "Val 1"
            table.cell(1, 1).text = "Val 2"

        path = _make_pptx(build)
        md = pptx_to_markdown(path)
        self.assertIn("| Col A | Col B |", md)
        self.assertIn("| --- | --- |", md)
        self.assertIn("| Val 1 | Val 2 |", md)

    def test_slide_separator(self):
        """Multiple slides should be separated by ---."""

        def build_slide1(prs):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Slide One"

        def build_slide2(prs):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Slide Two"

        path = _make_pptx(build_slide1, build_slide2)
        md = pptx_to_markdown(path)
        self.assertIn("---", md)
        self.assertIn("# Slide One", md)
        self.assertIn("# Slide Two", md)

    def test_speaker_notes(self):
        """Speaker notes should be rendered as a blockquote."""

        def build(prs):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = "Notes Slide"
            slide.notes_slide.notes_text_frame.text = "Important note here."

        path = _make_pptx(build)
        md = pptx_to_markdown(path)
        self.assertIn("> **Notes:**", md)
        self.assertIn("Important note here.", md)

    def test_inline_bold_in_body(self):
        """Bold runs inside a body shape should produce **bold** Markdown."""

        def build(prs):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Formatting"
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            para = tf.paragraphs[0]
            run = para.add_run()
            run.text = "bold text"
            run.font.bold = True

        path = _make_pptx(build)
        md = pptx_to_markdown(path)
        self.assertIn("**bold text**", md)

    def test_file_not_found(self):
        """Should raise FileNotFoundError for a missing file."""
        with self.assertRaises(FileNotFoundError):
            pptx_to_markdown(Path("/nonexistent/file.pptx"))

    def test_empty_presentation(self):
        """An empty presentation (no slides) should produce an empty string."""

        def build(prs):
            pass  # add no slides

        path = _make_pptx(build)
        md = pptx_to_markdown(path)
        self.assertEqual(md, "")


if __name__ == "__main__":
    unittest.main()
